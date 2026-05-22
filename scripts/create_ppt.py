# -*- coding: utf-8 -*-
"""
化学学术PPT创建脚本 - Chemistry Academic Presentation Builder

支持实验化学、理论计算化学、实验+理论混合三种论文类型的学术PPT生成。
Chemistry Academic Presentation Builder for experimental, computational, and hybrid papers.
"""
import sys
import os


def _safe_print(msg):
    """Windows-safe print that avoids GBK encoding crashes."""
    try:
        print(msg)
    except UnicodeEncodeError:
        print(msg.encode('ascii', errors='replace').decode('ascii'))


from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os.path


# ============================================================
# 配色主题 / Color Themes
# ============================================================

THEMES = {
    "academic": {
        "name": "学术经典",
        "bg":           (255, 255, 255),  # 白
        "title_color":  (0, 51, 102),      # 深蓝
        "text_color":   (51, 51, 51),      # 深灰
        "accent":       (180, 30, 30),      # 暗红
        "light_bg":     (240, 244, 248),   # 浅蓝灰
        "section_bg":   (0, 51, 102),      # 深蓝（章节页）
        "section_text": (255, 255, 255),   # 白
        "muted":        (140, 140, 140),   # 灰
        "border":       (220, 220, 220),   # 浅灰边框
        "table_header": (0, 51, 102),      # 深蓝
        "table_stripe": (240, 244, 248),   # 浅蓝灰
    },
    "molecular": {
        "name": "分子科技",
        "bg":           (248, 249, 250),
        "title_color":  (26, 82, 118),      # 钢蓝
        "text_color":   (44, 62, 80),
        "accent":       (231, 76, 60),      # 亮红
        "light_bg":     (235, 240, 245),
        "section_bg":   (26, 82, 118),
        "section_text": (255, 255, 255),
        "muted":        (149, 165, 166),
        "border":       (210, 218, 226),
        "table_header": (26, 82, 118),
        "table_stripe": (235, 240, 245),
    },
    "green": {
        "name": "绿色化学",
        "bg":           (247, 249, 244),
        "title_color":  (30, 86, 49),       # 深绿
        "text_color":   (51, 51, 51),
        "accent":       (212, 160, 23),     # 金
        "light_bg":     (238, 243, 233),
        "section_bg":   (30, 86, 49),
        "section_text": (255, 255, 255),
        "muted":        (150, 165, 140),
        "border":       (210, 220, 205),
        "table_header": (30, 86, 49),
        "table_stripe": (238, 243, 233),
    },
    "nature": {
        "name": "Nature 风格",
        "bg":           (255, 255, 255),
        "title_color":  (34, 34, 34),       # 近黑
        "text_color":   (68, 68, 68),
        "accent":       (0, 102, 204),      # 蓝
        "light_bg":     (248, 248, 248),
        "section_bg":   (34, 34, 34),
        "section_text": (255, 255, 255),
        "muted":        (160, 160, 160),
        "border":       (230, 230, 230),
        "table_header": (34, 34, 34),
        "table_stripe": (245, 245, 245),
    },
}


class ChemistryPPT:
    """化学学术PPT构建器 / Chemistry Academic PPT Builder"""

    def __init__(self, theme="academic"):
        if theme not in THEMES:
            raise ValueError(f"Unknown theme: {theme}. Choose from: {list(THEMES.keys())}")

        self.t = THEMES[theme]
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide_count = 0
        self._blank_layout = self.prs.slide_layouts[6]  # blank layout
        self._warnings = []
        self._errors = []
        self._missing_images = []
        self._slide_types = []
        self._theme_name = theme

    # ============================================================
    # 内部辅助 / Internal Helpers
    # ============================================================

    def _set_bg(self, slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color)

    def _add_textbox(self, slide, left, top, width, height,
                     text="", font_size=18, bold=False, color=None,
                     alignment=PP_ALIGN.LEFT, font_name=None, word_wrap=True):
        """添加文本框并返回 text_frame"""
        tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                       Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*(color or self.t["text_color"]))
        p.alignment = alignment
        if font_name:
            p.font.name = font_name
        return tf

    def _add_multi_paragraph(self, tf, texts, font_size=18, color=None,
                              bold=False, space_after=Pt(8)):
        """在已有 text_frame 中添加多个段落"""
        if color is None:
            color = self.t["text_color"]
        for i, text in enumerate(texts):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.color.rgb = RGBColor(*color)
            p.space_after = space_after
        return tf

    def _add_line(self, slide, x1, y1, x2, y2, color=None, width=Pt(1)):
        """添加装饰线"""
        connector = slide.shapes.add_connector(
            1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        connector.line.color.rgb = RGBColor(*(color or self.t["accent"]))
        connector.line.width = width
        return connector

    def _add_page_number(self, slide, num):
        """右下角页码"""
        self._add_textbox(
            slide, 12.0, 7.0, 1.2, 0.4,
            str(num), font_size=10, color=self.t["muted"],
            alignment=PP_ALIGN.RIGHT)

    def _add_subtitle_line(self, slide, text, top=1.0):
        """标题下方小字副标题"""
        self._add_textbox(
            slide, 0.7, top, 11.9, 0.4,
            text, font_size=12, color=self.t["muted"],
            alignment=PP_ALIGN.LEFT)

    # ============================================================
    # 幻灯片类型 / Slide Types
    # ============================================================

    def add_title_slide(self, title_cn, title_en="", authors="",
                         journal="", doi=""):
        """封面页 / Title Slide"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self.slide_count += 1
        self._slide_types.append("title")
        self._set_bg(slide, self.t["bg"])

        # 顶部装饰线
        self._add_line(slide, 0.7, 1.5, 12.6, 1.5, self.t["accent"], Pt(3))

        # 中文标题
        self._add_textbox(
            slide, 0.7, 1.8, 11.9, 2.0,
            title_cn, font_size=38, bold=True,
            color=self.t["title_color"], alignment=PP_ALIGN.LEFT)

        # 英文标题
        if title_en:
            self._add_textbox(
                slide, 0.7, 3.6, 11.9, 1.2,
                title_en, font_size=20, bold=False,
                color=self.t["muted"], alignment=PP_ALIGN.LEFT)

        # 作者
        if authors:
            self._add_textbox(
                slide, 0.7, 5.0, 11.9, 0.5,
                authors, font_size=16, color=self.t["text_color"],
                alignment=PP_ALIGN.LEFT)

        # 期刊信息 + DOI
        info_lines = []
        if journal:
            info_lines.append(journal)
        if doi:
            info_lines.append(f"DOI: {doi}")

        if info_lines:
            self._add_textbox(
                slide, 0.7, 5.5, 11.9, 0.5,
                "  |  ".join(info_lines), font_size=12,
                color=self.t["muted"], alignment=PP_ALIGN.LEFT)

        # 底部装饰条
        self._add_line(slide, 0.7, 6.5, 12.6, 6.5, self.t["border"], Pt(0.5))

        self._add_page_number(slide, self.slide_count)

    def add_section_slide(self, title, subtitle=""):
        """章节分隔页 / Section Divider"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self.slide_count += 1
        self._slide_types.append("section")
        self._set_bg(slide, self.t["section_bg"])

        # 左装饰线
        self._add_line(slide, 0.8, 2.6, 0.8, 5.0, self.t["accent"], Pt(3))

        self._add_textbox(
            slide, 1.5, 3.0, 10.5, 1.5,
            title, font_size=40, bold=True,
            color=self.t["section_text"], alignment=PP_ALIGN.LEFT)

        if subtitle:
            self._add_textbox(
                slide, 1.5, 4.3, 10.5, 0.6,
                subtitle, font_size=16,
                color=tuple(min(c + 80, 255) for c in self.t["section_bg"]),
                alignment=PP_ALIGN.LEFT)

        self._add_page_number(slide, self.slide_count)

    def add_content_slide(self, title, bullets, subtitle="", notes=""):
        """文字要点页 / Content Slide"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self.slide_count += 1
        self._slide_types.append("content")
        self._set_bg(slide, self.t["bg"])

        # 标题
        self._add_textbox(
            slide, 0.7, 0.4, 11.9, 0.8,
            title, font_size=32, bold=True, color=self.t["title_color"])

        # 标题下细线
        self._add_line(slide, 0.7, 1.15, 8.0, 1.15, self.t["accent"], Pt(1.5))

        if subtitle:
            self._add_subtitle_line(slide, subtitle, top=1.25)

        # 正文要点
        top = 1.6 if not subtitle else 1.7
        tb = slide.shapes.add_textbox(
            Inches(0.7), Inches(top), Inches(11.6), Inches(5.3))
        tf = tb.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(*self.t["text_color"])
            p.space_after = Pt(14)
            p.level = 0

        # 备注
        if notes:
            self._add_textbox(
                slide, 0.7, 6.8, 11.9, 0.4,
                f"  {notes}", font_size=10, color=self.t["muted"])

        self._add_page_number(slide, self.slide_count)
        return slide

    def add_figure_slide(self, title, figure_path, bullets=None,
                          figure_label="", caption="",
                          layout="figure_right", notes=""):
        """图表+说明页 / Figure + Explanation Slide

        Args:
            layout: "figure_right" (图右文左), "figure_top" (图上文下),
                    "figure_full" (全图+底部说明), "figure_left" (图左文右)
        """
        slide = self.prs.slides.add_slide(self._blank_layout)
        self.slide_count += 1
        self._slide_types.append("figure")
        self._set_bg(slide, self.t["bg"])

        # 标题
        self._add_textbox(
            slide, 0.7, 0.3, 11.9, 0.7,
            title, font_size=28, bold=True, color=self.t["title_color"])
        self._add_line(slide, 0.7, 1.0, 7.0, 1.0, self.t["accent"], Pt(1.5))

        bullets = bullets or []
        img_ok = os.path.exists(figure_path) if figure_path else False
        if figure_path and not img_ok:
            self._missing_images.append(figure_path)

        if layout == "figure_right":
            fig_left = 7.2
            fig_width = 5.5
            text_width = 5.8

            if bullets:
                tb = slide.shapes.add_textbox(
                    Inches(0.7), Inches(1.3), Inches(text_width), Inches(5.3))
                tf = tb.text_frame
                tf.word_wrap = True
                for i, b in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = b
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(*self.t["text_color"])
                    p.space_after = Pt(10)

            if img_ok:
                try:
                    slide.shapes.add_picture(
                        figure_path, Inches(fig_left), Inches(1.3),
                        width=Inches(fig_width))
                except Exception as e:
                    self._errors.append(f"Failed to insert image '{figure_path}': {e}")
                    self._add_textbox(
                        slide, fig_left, 2.5, fig_width, 1.0,
                        f"[Image load error: {os.path.basename(figure_path)}]",
                        font_size=14, color=self.t["muted"],
                        alignment=PP_ALIGN.CENTER)
            elif figure_path:
                self._add_textbox(
                    slide, fig_left, 2.5, fig_width, 1.0,
                    f"[Figure not found: {os.path.basename(figure_path)}]",
                    font_size=14, color=self.t["muted"],
                    alignment=PP_ALIGN.CENTER)

        elif layout == "figure_top":
            if img_ok:
                try:
                    slide.shapes.add_picture(
                        figure_path, Inches(0.7), Inches(1.2),
                        height=Inches(3.2))
                except Exception as e:
                    self._errors.append(f"Failed to insert image '{figure_path}': {e}")

            if bullets:
                tb = slide.shapes.add_textbox(
                    Inches(0.7), Inches(4.6), Inches(11.6), Inches(2.5))
                tf = tb.text_frame
                tf.word_wrap = True
                for i, b in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = b
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(*self.t["text_color"])
                    p.space_after = Pt(8)

        elif layout == "figure_left":
            if img_ok:
                try:
                    slide.shapes.add_picture(
                        figure_path, Inches(0.5), Inches(1.3),
                        width=Inches(5.5))
                except Exception as e:
                    self._errors.append(f"Failed to insert image '{figure_path}': {e}")

            if bullets:
                tb = slide.shapes.add_textbox(
                    Inches(6.5), Inches(1.3), Inches(6.2), Inches(5.3))
                tf = tb.text_frame
                tf.word_wrap = True
                for i, b in enumerate(bullets):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = b
                    p.font.size = Pt(16)
                    p.font.color.rgb = RGBColor(*self.t["text_color"])
                    p.space_after = Pt(10)

        elif layout == "figure_full":
            if img_ok:
                try:
                    slide.shapes.add_picture(
                        figure_path, Inches(0.5), Inches(1.2),
                        width=Inches(12.3))
                except Exception as e:
                    self._errors.append(f"Failed to insert image '{figure_path}': {e}")

            if bullets:
                y = 6.0
                for b in bullets:
                    self._add_textbox(
                        slide, 0.7, y, 11.9, 0.4,
                        b, font_size=12, color=self.t["text_color"])
                    y += 0.3

        # 图注/来源
        if figure_label or caption:
            label_text = figure_label or ""
            if caption:
                label_text += f"  |  {caption}" if label_text else caption
            self._add_textbox(
                slide, 0.7, 6.8, 11.9, 0.4,
                label_text, font_size=9, color=self.t["muted"])

        if notes:
            self._add_textbox(
                slide, 0.7, 6.55, 11.9, 0.3,
                f"  {notes}", font_size=9, color=self.t["muted"])

        self._add_page_number(slide, self.slide_count)
        return slide

    def add_table_slide(self, title, headers, rows, notes=""):
        """数据表格页 / Table Slide"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self.slide_count += 1
        self._slide_types.append("table")
        self._set_bg(slide, self.t["bg"])

        self._add_textbox(
            slide, 0.7, 0.3, 11.9, 0.7,
            title, font_size=28, bold=True, color=self.t["title_color"])
        self._add_line(slide, 0.7, 1.0, 7.0, 1.0, self.t["accent"], Pt(1.5))

        n_rows = len(rows) + 1  # +1 for header
        n_cols = len(headers)

        # 表格定位
        table_left = Inches(0.7)
        table_top = Inches(1.4)
        table_width = Inches(12.0)
        row_height = Inches(min(0.45, 4.5 / max(n_rows, 1)))
        table_height = row_height * n_rows

        table_shape = slide.shapes.add_table(
            n_rows, n_cols, table_left, table_top,
            table_width, table_height)
        table = table_shape.table

        # 设置列宽（等分）
        col_width = int(Emu(table_width) / n_cols)
        for col_idx in range(n_cols):
            table.columns[col_idx].width = col_width

        # 表头
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(14)
                paragraph.font.bold = True
                paragraph.font.color.rgb = RGBColor(255, 255, 255)
                paragraph.alignment = PP_ALIGN.CENTER
            # 表头背景
            tcPr = cell._tc.get_or_add_tcPr()
            solidFill = cell._tc.makeelement(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill', {})
            srgbClr = cell._tc.makeelement(
                '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr',
                {'val': '{:02X}{:02X}{:02X}'.format(*self.t["table_header"])})
            solidFill.append(srgbClr)
            tcPr.append(solidFill)

        # 数据行
        for row_idx, row in enumerate(rows):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(13)
                    paragraph.font.color.rgb = RGBColor(*self.t["text_color"])
                    paragraph.alignment = PP_ALIGN.CENTER
                # 斑马纹
                if row_idx % 2 == 0:
                    tcPr = cell._tc.get_or_add_tcPr()
                    solidFill = cell._tc.makeelement(
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill', {})
                    srgbClr = cell._tc.makeelement(
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr',
                        {'val': '{:02X}{:02X}{:02X}'.format(*self.t["table_stripe"])})
                    solidFill.append(srgbClr)
                    tcPr.append(solidFill)

        if notes:
            self._add_textbox(
                slide, 0.7, 6.8, 11.9, 0.4,
                f"  {notes}", font_size=10, color=self.t["muted"])

        self._add_page_number(slide, self.slide_count)
        return slide

    def add_summary_slide(self, title, bullets, notes=""):
        """总结/结论页 / Summary Slide"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self.slide_count += 1
        self._slide_types.append("summary")
        self._set_bg(slide, self.t["light_bg"])

        # 顶部色条
        self._add_line(slide, 0, 0, 13.333, 0, self.t["title_color"], Pt(4))

        self._add_textbox(
            slide, 0.7, 0.6, 11.9, 0.8,
            title, font_size=34, bold=True, color=self.t["title_color"])

        # 要点
        tb = slide.shapes.add_textbox(
            Inches(0.7), Inches(1.7), Inches(11.6), Inches(5.0))
        tf = tb.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(*self.t["text_color"])
            p.space_after = Pt(16)

        if notes:
            self._add_textbox(
                slide, 0.7, 6.8, 11.9, 0.4,
                f"  {notes}", font_size=10, color=self.t["muted"])

        self._add_page_number(slide, self.slide_count)
        return slide

    def add_thankyou_slide(self, title="谢谢！欢迎提问", subtitle="Thank you & Questions"):
        """致谢页 / Thank You Slide"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self.slide_count += 1
        self._slide_types.append("thankyou")
        self._set_bg(slide, self.t["section_bg"])

        self._add_textbox(
            slide, 0, 2.5, 13.333, 1.5,
            title, font_size=48, bold=True,
            color=self.t["section_text"], alignment=PP_ALIGN.CENTER)

        if subtitle:
            self._add_textbox(
                slide, 0, 4.2, 13.333, 0.8,
                subtitle, font_size=20,
                color=self.t["muted"], alignment=PP_ALIGN.CENTER)

        self._add_page_number(slide, self.slide_count)
        return slide

    def add_image_grid_slide(self, title, image_items, cols=2, notes=""):
        """多图网格页 / Multi-Image Grid Slide

        Args:
            image_items: list of (image_path, caption) tuples
            cols: 每行图片数
        """
        slide = self.prs.slides.add_slide(self._blank_layout)
        self.slide_count += 1
        self._slide_types.append("image_grid")
        self._set_bg(slide, self.t["bg"])

        self._add_textbox(
            slide, 0.7, 0.3, 11.9, 0.7,
            title, font_size=28, bold=True, color=self.t["title_color"])
        self._add_line(slide, 0.7, 1.0, 7.0, 1.0, self.t["accent"], Pt(1.5))

        n_items = len(image_items)
        rows = (n_items + cols - 1) // cols

        cell_w = 11.8 / cols
        cell_h = 5.2 / rows

        for idx, (img_path, caption) in enumerate(image_items):
            r = idx // cols
            c = idx % cols
            left = 0.7 + c * cell_w + 0.1
            top = 1.3 + r * cell_h + 0.1

            if os.path.exists(img_path):
                try:
                    slide.shapes.add_picture(
                        img_path, Inches(left), Inches(top),
                        width=Inches(cell_w - 0.2), height=Inches(cell_h - 0.5))
                except Exception:
                    pass

            if caption:
                self._add_textbox(
                    slide, left, top + cell_h - 0.5, cell_w - 0.2, 0.4,
                    caption, font_size=9, color=self.t["muted"],
                    alignment=PP_ALIGN.CENTER)

        if notes:
            self._add_textbox(
                slide, 0.7, 6.8, 11.9, 0.4,
                f"  {notes}", font_size=10, color=self.t["muted"])

        self._add_page_number(slide, self.slide_count)
        return slide

    # ============================================================
    # 保存 / Save
    # ============================================================

    def save(self, output_path):
        """保存 PPTX 文件"""
        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        self.prs.save(output_path)
        _safe_print(f"[ChemistryPPT] Saved: {output_path}")
        _safe_print(f"[ChemistryPPT] Slides: {self.slide_count}  |  Theme: {self.t['name']}")

        type_counts = {}
        for st in self._slide_types:
            type_counts[st] = type_counts.get(st, 0) + 1
        type_summary = ", ".join(f"{v}x {k}" for k, v in type_counts.items())
        _safe_print(f"[ChemistryPPT] Types: {type_summary}")

        if self._missing_images:
            _safe_print(f"[ChemistryPPT] Missing images: {len(self._missing_images)}")
            for m in self._missing_images:
                print(f"  - {m}")
        if self._warnings:
            _safe_print(f"[ChemistryPPT] Warnings: {len(self._warnings)}")
            for w in self._warnings:
                print(f"  - {w}")
        if self._errors:
            _safe_print(f"[ChemistryPPT] Errors: {len(self._errors)}")
            for e in self._errors:
                print(f"  - {e}")

        return output_path

    def get_report(self):
        """获取构建报告"""
        return {
            "theme": self._theme_name,
            "theme_display": self.t["name"],
            "total_slides": self.slide_count,
            "slide_types": self._slide_types,
            "missing_images": self._missing_images,
            "warnings": self._warnings,
            "errors": self._errors,
        }

    def save_report(self, output_path):
        """保存 JSON 构建报告"""
        import json
        report = self.get_report()
        report_path = output_path.replace('.pptx', '_report.json')
        with open(report_path, 'w', encoding='utf-8') as f:
            json.dump(report, f, indent=2, ensure_ascii=False)
        _safe_print(f"[ChemistryPPT] Report saved: {report_path}")
        return report_path


# ============================================================
# 命令行入口 / CLI Entry
# ============================================================

def main():
    print("=" * 50)
    print("Chemistry Academic PPT Builder")
    print("=" * 50)
    print()
    print("Usage (Python API):")
    print("  from create_ppt import ChemistryPPT")
    print("  ppt = ChemistryPPT(theme='academic')")
    print("  ppt.add_title_slide(...)")
    print("  ppt.save('output.pptx')")
    print()
    print("Available themes:", ", ".join(THEMES.keys()))
    print()
    print("For automated paper-to-PPT conversion, use this skill's")
    print("SKILL.md workflow: read paper → classify → build slides.")


if __name__ == "__main__":
    main()
